

import streamlit as st


import openai
import os
from dotenv import load_dotenv
import PyPDF2
import docx
from pptx import Presentation

load_dotenv()
openai_api_key = os.getenv("OPENAI_API_KEY")


st.set_page_config(page_title="Sales Agent Prototype", layout="centered")
st.title(":bar_chart: Sales Agent Prototype")

# Sidebar for instructions
with st.sidebar:
    st.header("Instructions")
    st.markdown("""
    1. Fill in all required fields.
    2. (Optional) Upload a product overview document.
    3. Click **Generate Insights** to receive a one-page summary.
    4. Review the output and use it for your sales preparation.
    """)

st.write("""
<span style='font-size:18px;'>Enter the details below to generate account insights for your sales opportunity.</span>
""", unsafe_allow_html=True)


# Input fields
product_name = st.text_input("Product Name", help="What product are you selling?")
company_url = st.text_input("Company URL", help="URL of the company you are targeting")
product_category = st.text_input("Product Category", help="E.g., Data Warehousing, Cloud Data Platform")
competitors = st.text_area("Competitors (URLs)", help="Enter one URL per line")
value_proposition = st.text_input("Value Proposition", help="Summarize the product's value in one sentence")
target_customer = st.text_input("Target Customer", help="Name of the person you are trying to sell to")
manual_leaders = st.text_area("Known Company Leaders (optional, one per line)", help="Enter names of key leaders at the prospect company")
uploaded_file = st.file_uploader("Optional: Upload a product overview sheet or deck", type=["pdf", "docx", "pptx", "txt"])

# Customizable output template: section selection
st.markdown("**Select which sections to include in the one-pager:**")
section_options = {
    "Company Strategy": "company_strategy",
    "Competitor Mentions": "competitor_mentions",
    "Leadership Information": "leadership_info",
    "Product/Strategy Summary": "product_strategy"
}
selected_sections = st.multiselect(
    "Sections",
    options=list(section_options.keys()),
    default=list(section_options.keys()),
    help="Choose which sections to include in the generated insights."
)


def extract_text_from_file(uploaded_file):
    if uploaded_file is None:
        return ""
    if uploaded_file.type == "application/pdf":
        try:
            reader = PyPDF2.PdfReader(uploaded_file)
            text = " ".join(page.extract_text() or "" for page in reader.pages)
            return text[:2000]  # Limit to 2000 chars
        except Exception:
            return "[Could not extract text from PDF]"
    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
        try:
            doc = docx.Document(uploaded_file)
            text = " ".join([p.text for p in doc.paragraphs])
            return text[:2000]
        except Exception:
            return "[Could not extract text from DOCX]"
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        try:
            prs = Presentation(uploaded_file)
            text = " ".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            return text[:2000]
        except Exception:
            return "[Could not extract text from PPTX]"
    elif uploaded_file.type == "text/plain":
        try:
            return uploaded_file.read().decode("utf-8")[:2000]
        except Exception:
            return "[Could not extract text from TXT]"
    else:
        return "[Unsupported file type]"


# Multi-agent prompt builder for each section
def build_section_prompt(section_key, doc_text, manual_leaders=None):
    base = f"""
Product Name: {product_name}
Company URL: {company_url}
Product Category: {product_category}
Competitors: {competitors}
Value Proposition: {value_proposition}
Target Customer: {target_customer}
Product Overview Document Extract (if provided):
{doc_text}
"""
    if manual_leaders and section_key == "leadership_info":
        base += f"\nKnown Company Leaders (user provided):\n{manual_leaders}\n"
    if section_key == "company_strategy":
        return (
            "You are a professional sales insights assistant. "
            "Summarize the company's activities and direction in the relevant industry. "
            "Reference any public statements, press releases, or articles by key executives. "
            "Mention relevant job postings or technology stack indicators if available.\n" + base
        )
    elif section_key == "competitor_mentions":
        return (
            "You are a professional sales insights assistant. "
            "Note any public information about the listed competitors and their relationship to the target company.\n" + base
        )
    elif section_key == "leadership_info":
        return (
            "You are a professional sales insights assistant. "
            "List key leaders at the prospect company, especially those quoted in recent press releases or articles.\n" + base
        )
    elif section_key == "product_strategy":
        return (
            "You are a professional sales insights assistant. "
            "For public companies, include insights from annual reports or other relevant documents.\n" + base
        )
    else:
        return base

if st.button("Generate Insights"):
    if not openai_api_key:
        st.error("OpenAI API key not found. Please set it in your .env file.")
    elif not product_name or not company_url:
        st.warning("Please fill in at least Product Name and Company URL.")
    elif not selected_sections:
        st.warning("Please select at least one section to include in the output.")
    else:
        with st.spinner("Generating insights with GPT (multi-agent)..."):
            try:
                client = openai.OpenAI(api_key=openai_api_key)
                doc_text = extract_text_from_file(uploaded_file)
                # Multi-agent: run LLM for each selected section
                sections = [(title, section_options[title]) for title in selected_sections]
                results = {}
                for title, key in sections:
                    prompt = build_section_prompt(key, doc_text, manual_leaders)
                    response = client.chat.completions.create(
                        model="gpt-3.5-turbo",
                        messages=[
                            {"role": "system", "content": "You are a helpful sales assistant agent."},
                            {"role": "user", "content": prompt}
                        ],
                        max_tokens=500,
                        temperature=0.7
                    )
                    results[title] = response.choices[0].message.content.strip()
                st.success("Insights generated!")
                st.markdown("---")
                st.markdown("<h3>Sales Account Insights</h3>", unsafe_allow_html=True)
                # Section summaries/tooltips
                section_summaries = {
                    "Company Strategy": "Summary of the company's activities, direction, and public statements in the relevant industry.",
                    "Competitor Mentions": "Mentions and analysis of competitors relevant to the target company.",
                    "Leadership Information": "Key leaders at the prospect company, especially those quoted in recent press releases or articles.",
                    "Product/Strategy Summary": "Insights from annual reports or other relevant documents about the company's product or strategy."
                }
                for title, _ in sections:
                    st.markdown(f"### {title}")
                    if title in section_summaries:
                        st.caption(section_summaries[title])
                    st.markdown(results[title], unsafe_allow_html=True)
                # Data Visualization: Competitor Mentions Bar Chart
                import matplotlib.pyplot as plt
                import pandas as pd
                if "Competitor Mentions" in results and competitors.strip():
                    competitor_list = [c.strip() for c in competitors.splitlines() if c.strip()]
                    competitor_counts = {}
                    competitor_text = results["Competitor Mentions"].lower()
                    for comp in competitor_list:
                        count = competitor_text.count(comp.lower())
                        competitor_counts[comp] = count
                    if competitor_counts:
                        st.markdown("#### Competitor Mention Frequency (in Insights)")
                        df = pd.DataFrame({"Competitor": list(competitor_counts.keys()), "Mentions": list(competitor_counts.values())})
                        fig, ax = plt.subplots()
                        df.plot(kind="bar", x="Competitor", y="Mentions", legend=False, ax=ax, color="#1f77b4")
                        ax.set_ylabel("Mentions")
                        ax.set_xlabel("")
                        ax.set_title("")
                        st.pyplot(fig)
                st.markdown("---")
            except openai.APIStatusError as api_err:
                if api_err.status_code == 429:
                    st.error("You have exceeded your OpenAI API quota or rate limit. Please check your usage and billing details.")
                elif api_err.status_code == 401:
                    st.error("Invalid OpenAI API key. Please check your .env file.")
                else:
                    st.error(f"OpenAI API error: {api_err}")
            except Exception as e:
                st.error(f"Unexpected error generating insights: {e}")
